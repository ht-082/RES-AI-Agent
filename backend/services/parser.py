import os
import re
import fitz  # PyMuPDF
import docx
import openpyxl
import logging
from collections import Counter
from services.vision_parser import parse_with_vision_api

logger = logging.getLogger(__name__)


# ── 매직바이트 게이트 (전략 §3.0) ────────────────────────────────────
# 확장자를 믿지 않는다. 실측: ole2가 .docx로 위장 8건(전부 DRM), 그룹웨어 래핑 pdf 1건.

class ParserGateError(Exception):
    """파싱 불가 사유를 담아 호출측(tasks)이 fail_reason으로 기록하게 한다."""

    def __init__(self, reason):
        super().__init__(reason)
        self.reason = reason


_MAGIC_SIGS = [(b'PK\x03\x04', 'zip'), (b'\xd0\xcf\x11\xe0', 'ole2'),
               (b'%PDF', 'pdf'), (b'{\\rtf', 'rtf')]
_MAGIC_EXPECT = {'pdf': 'pdf', 'docx': 'zip', 'doc': 'ole2', 'xlsx': 'zip', 'xlsm': 'zip',
                 'pptx': 'zip', 'hwpx': 'zip', 'hwp': 'ole2'}
# DRM/암호화 컨테이너 마커 (ASCII·UTF-16 모두 확인 — Fasoo 'FSD', MS 암호화 래퍼)
_DRM_MARKERS = [b'FSD', b'EncryptedPackage', b'EncryptionInfo',
                'EncryptedPackage'.encode('utf-16-le'), 'Encrypt'.encode('utf-16-le')]


def _sniff_magic(file_path):
    with open(file_path, 'rb') as f:
        head = f.read(8)
    for sig, name in _MAGIC_SIGS:
        if head.startswith(sig):
            return name
    return 'unknown'


def gate_check(file_path, ext):
    """파일 서명 검증. 정상이면 옵션 dict, 불가면 ParserGateError.

    반환 옵션: {'pdf_offset': N} — 그룹웨어(전자결재) 래핑 PDF의 실제 %PDF 시작 위치.
    """
    magic = _sniff_magic(file_path)
    expect = _MAGIC_EXPECT.get(ext)
    if expect is None or magic == expect:
        return {}

    # pdf 확장자인데 서명이 다르면: 파일 안에 %PDF가 통째로 들어있는지 탐색
    # (실측: 핸디소프트 전자결재 헤더 429바이트 뒤 정상 PDF)
    if ext == 'pdf':
        with open(file_path, 'rb') as f:
            head = f.read(65536)
        off = head.find(b'%PDF')
        if off > 0:
            logger.info(f"래핑된 PDF 감지: 오프셋 {off}부터 읽음 ({os.path.basename(file_path)})")
            return {'pdf_offset': off}

    # .hwpx인데 실체가 ole2면 그냥 hwp다(사용자가 확장자만 바꿔 저장한 경우).
    # 내용은 멀쩡하므로 차단하지 않고 hwp 파서로 넘긴다.
    if ext == 'hwpx' and magic == 'ole2':
        logger.info(f"hwpx 확장자이나 실체는 hwp — hwp 파서로 처리 ({os.path.basename(file_path)})")
        return {'as_ext': 'hwp'}

    # ole2 위장이면 DRM 여부를 구분해 기록 (권한 문제 vs 포맷 문제)
    if magic == 'ole2':
        size = os.path.getsize(file_path)
        with open(file_path, 'rb') as f:
            data = f.read(min(size, 4_000_000))
        if any(m in data for m in _DRM_MARKERS):
            raise ParserGateError('drm-protected')

        # DRM이 아닌 ole2 = 그냥 구형 포맷이다. 확장자만 .docx로 바뀐 .doc 파일이
        # 코퍼스에 4건 있었고 전부 차단돼 있었다(실측: antiword로 1.4만~17만자 추출).
        # DRM 검사 **뒤에** 두는 것이 중요하다 — 순서를 바꾸면 DRM 문서가 이리 샌다.
        if ext in ('docx', 'doc'):
            logger.info(f"docx 확장자이나 실체는 레거시 doc — doc 파서로 처리 "
                        f"({os.path.basename(file_path)})")
            return {'as_ext': 'doc'}

    raise ParserGateError(f'format-mismatch({magic})')


def _rows_to_markdown(rows, group_size=20):
    """표 행렬 → Markdown 표 문자열 목록 (전략 §3.6).

    대형 표는 헤더를 반복하며 group_size 행씩 분할 — 어느 청크든 자기설명적이 되게.
    헤더+데이터 2행 미만이면 표로 보지 않고 빈 목록을 반환한다(본문 텍스트로 남김).
    """
    clean = [[('' if c is None else str(c).replace('|', '/').replace('\n', ' ').strip())
              for c in r] for r in rows]
    clean = [r for r in clean if any(r)]
    if len(clean) < 2:
        return []
    header, data = clean[0], clean[1:]
    head_md = '| ' + ' | '.join(header) + ' |'
    sep_md = '|' + '---|' * len(header)
    out = []
    for i in range(0, len(data), group_size):
        lines = [head_md, sep_md]
        for r in data[i:i + group_size]:
            r = (r + [''] * len(header))[:len(header)]
            lines.append('| ' + ' | '.join(r) + ' |')
        out.append('\n'.join(lines))
    return out

# ── 머리말/꼬리말 제거 ────────────────────────────────────────────────
# 법제처 PDF는 모든 페이지 상단에 법령명, 하단에 쪽번호를 반복해 넣는다.
# 그대로 두면 청크마다 같은 문자열이 박혀 검색·리랭킹을 오염시킨다.
# (실측: 「전기사업법」42/42p에 '전기사업법', 「개발행위허가」28/30p에 'p.N')
EDGE_LINES = 2        # 페이지 위/아래 몇 줄까지를 머리말·꼬리말 후보로 볼지
EDGE_MIN_RATIO = 0.6  # 전체 페이지의 몇 할 이상에서 반복되면 제거할지
EDGE_MAX_LEN = 60     # 이보다 긴 줄은 본문으로 간주해 건드리지 않음


def _normalize_line(line):
    """쪽번호처럼 숫자만 달라지는 줄을 하나로 묶기 위해 숫자를 #로 치환"""
    return re.sub(r'\d+', '#', line.strip())


def strip_repeated_edges(pages):
    """페이지 가장자리(첫/마지막 줄)에 반복되는 머리말·꼬리말을 제거한다.

    위치를 따지는 이유: 법령의 '[전문개정 2020. 6. 9.]'처럼 본문 중간에 정당하게
    반복되는 문자열이 있다. 단순 빈도로 지우면 그런 내용까지 날아간다.
    """
    if len(pages) < 3:  # 표본이 적으면 반복 여부를 신뢰할 수 없다
        return pages

    per_page = []
    head_counter, tail_counter = Counter(), Counter()
    for page in pages:
        lines = [l for l in (page.get('text') or '').split('\n') if l.strip()]
        per_page.append(lines)
        for line in lines[:EDGE_LINES]:
            if len(line.strip()) <= EDGE_MAX_LEN:
                head_counter[_normalize_line(line)] += 1
        for line in lines[-EDGE_LINES:]:
            if len(line.strip()) <= EDGE_MAX_LEN:
                tail_counter[_normalize_line(line)] += 1

    threshold = max(3, int(len(pages) * EDGE_MIN_RATIO))
    drop = {t for t, c in head_counter.items() if c >= threshold}
    drop |= {t for t, c in tail_counter.items() if c >= threshold}
    if not drop:
        return pages

    cleaned, removed = [], 0
    for page, lines in zip(pages, per_page):
        total = len(lines)
        kept = []
        for idx, line in enumerate(lines):
            at_edge = idx < EDGE_LINES or idx >= total - EDGE_LINES
            if at_edge and _normalize_line(line) in drop:
                removed += 1
                continue
            kept.append(line)
        item = dict(page)
        item['text'] = '\n'.join(kept).strip()
        cleaned.append(item)

    logger.info(f"머리말/꼬리말 {len(drop)}종 · {removed}줄 제거 (페이지 {len(pages)})")
    # 제거 후 내용이 비어버린 페이지는 청킹 대상에서 뺀다(기존 파서 동작과 동일)
    return [p for p in cleaned if (p.get('text') or '').strip()]

def parse_pdf(file_path, pdf_offset=0):
    """PDF 파싱 (페이지 단위) + 표 구조 보존 (전략 §3.1, §3.6)

    - pdf_offset: 그룹웨어 래핑 파일의 %PDF 시작 위치 (gate_check가 알려줌)
    - 표는 fitz.find_tables로 감지, Markdown 직렬화에 성공한 표만 본문에서
      제거(redaction)한다 — 실패 시 본문 텍스트에 그대로 남아 내용 손실이 없다.
    - 스캔본(추출 50자 미만)은 Vision OCR로 우회하고 각 항목에 ocr=True를 남긴다.
    """
    pages, tables = [], []
    try:
        if pdf_offset:
            with open(file_path, 'rb') as f:
                data = f.read()[pdf_offset:]
            doc = fitz.open(stream=data, filetype='pdf')
        else:
            doc = fitz.open(file_path)

        for page_idx, page in enumerate(doc):
            # 1) 표 감지 → 직렬화 성공분만 본문에서 제거
            page_tables = []
            try:
                for tab in page.find_tables().tables:
                    md_parts = _rows_to_markdown(tab.extract())
                    if md_parts:
                        page_tables.append((tab.bbox, md_parts))
                if page_tables:
                    for bbox, _ in page_tables:
                        page.add_redact_annot(bbox)
                    page.apply_redactions()
            except Exception as e:
                logger.warning(f"표 추출 실패(p.{page_idx + 1}, {os.path.basename(file_path)}): {e}")
                page_tables = []

            # 2) (표 영역이 제거된) 본문 텍스트
            text = page.get_text().strip()
            if text:
                pages.append({
                    'text': text,
                    'page_number': page_idx + 1,
                    'section_title': f"페이지 {page_idx + 1}",
                    'sheet_name': '',
                    'cell_range': ''
                })
            for _, md_parts in page_tables:
                for md in md_parts:
                    tables.append({
                        'text': md,
                        'page_number': page_idx + 1,
                        'section_title': f"표 (페이지 {page_idx + 1})",
                        'kind': 'table',
                        'sheet_name': '',
                        'cell_range': ''
                    })
        doc.close()
    except Exception as e:
        logger.error(f"PDF 파싱 실패 ({file_path}): {e}")

    # 스캔본 판별 → Vision OCR (전략 §3.1: 전량 OCR + ocr 마킹)
    total_length = sum(len(p.get('text', '')) for p in pages)
    if total_length < 50 and not tables:
        logger.warning(f"PDF 텍스트 추출량 미달({total_length}자). Vision OCR로 우회: {file_path}")
        ocr_items = parse_with_vision_api(file_path)
        for item in ocr_items:
            item['ocr'] = True
        return ocr_items

    return pages + tables

def parse_docx(file_path):
    """
    Word 파일 파싱 (단락 단위 -> 적정 묶음으로 분할)
    """
    paragraphs_data = []
    try:
        doc = docx.Document(file_path)
        # 단락별 텍스트를 모으고, 적정 길이(예: 3단락) 단위로 페이지화하듯이 묶음
        current_text = []
        para_count = 0
        section_idx = 1
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            current_text.append(text)
            para_count += 1
            
            # 5개 단락마다 또는 800자가 넘어가면 하나의 섹션으로 묶음
            joined_text = "\n".join(current_text)
            if para_count >= 5 or len(joined_text) >= 800:
                paragraphs_data.append({
                    'text': joined_text,
                    'page_number': section_idx,
                    'section_title': f"섹션 {section_idx}",
                    'sheet_name': '',
                    'cell_range': ''
                })
                current_text = []
                para_count = 0
                section_idx += 1
                
        # 남은 텍스트 처리
        if current_text:
            paragraphs_data.append({
                'text': "\n".join(current_text),
                'page_number': section_idx,
                'section_title': f"섹션 {section_idx}",
                'sheet_name': '',
                'cell_range': ''
            })
            
        # 표(Table) 파싱 — Markdown 직렬화 (전략 §3.6). 위치 보존은 한계 인정(문서 말미 배치),
        # page 번호는 본문 섹션과 겹치지 않게 마지막 섹션 다음부터 부여한다(기존 충돌 버그 수정).
        table_page = section_idx + 1
        for table_idx, table in enumerate(doc.tables):
            rows = [[cell.text for cell in row.cells] for row in table.rows]
            md_parts = _rows_to_markdown(rows)
            if md_parts:
                for part_idx, md in enumerate(md_parts):
                    paragraphs_data.append({
                        'text': f"[표 {table_idx + 1}]\n{md}",
                        'page_number': table_page,
                        'section_title': f"표 {table_idx + 1}",
                        'kind': 'table',
                        'sheet_name': '',
                        'cell_range': ''
                    })
                    table_page += 1
            else:
                # 1행짜리 등 표로 볼 수 없는 것은 기존 방식으로 텍스트 보존
                flat = [" | ".join(c.strip() for c in r if c.strip()) for r in rows]
                flat = [f for f in flat if f]
                if flat:
                    paragraphs_data.append({
                        'text': f"[표 {table_idx + 1}]\n" + "\n".join(flat),
                        'page_number': table_page,
                        'section_title': f"표 {table_idx + 1}",
                        'sheet_name': '',
                        'cell_range': ''
                    })
                    table_page += 1
    except Exception as e:
        logger.error(f"Docx 파싱 실패 ({file_path}): {e}")
    return paragraphs_data

# ── Excel 부피 제어 (재무모델 대응) ──────────────────────────────────
# 실측(재무모델 xlsm 8건): 원시 텍스트는 파일당 54만자인데 행 서술화('헤더: 값' 쌍)를
# 거치면 265만자로 **4.9배** 부푼다. 분기 시계열 시트는 열이 수십 개라 헤더 문자열이
# 그만큼 반복되기 때문이다. 8건 합계 1,042만자 = 청크 약 1만개로, v1.0 코퍼스 전체
# (10,777청크)보다 크다. 행 서술화는 열이 적을 때만 이득이다.
XLSX_PAIR_MAX_COLS = 12

# 재무모델은 요약·가정 시트만 인덱싱한다.
# 근거: 8개 파일 170여 시트를 훑어보니 분기 캐시플로·재무제표 격자는 숫자셀이 84~97%다.
# 사람은 챗봇에 "IRR 얼마야"·"사업비 총액"을 묻지, 2024년 3분기 원리금을 묻지 않는다
# (그건 엑셀을 연다). 상세 격자는 비용만 크고 검색에 기여하지 않는다.
FM_SHEET_WHITELIST = re.compile(
    r'report|summary|result|보고서|요약'
    r'|irr|수익률|valuation|lcoe|경제성'
    r'|assum(?!_book)|가정|input|전제'          # Assum_Book(가정 상세)은 제외
    r'|tic|tpc|투자비|총사업비|funding|재원조달'
    r'|revenue|^rev$|발전량|opex|capex'
    r'|sensitivity|민감도',
    re.I)

# 화이트리스트를 통과해도 제외하는 것: 이름 끝의 주기 접미사.
# 8개 파일 공통 관례로 (Q)=분기·(Y)=연간·(FY)=회계연도 시계열 격자를 뜻한다.
# 'IRR'은 요약이지만 'IRR(Q)'는 264행짜리 분기 격자(8.9만자)다 — 같은 단어라도
# 접미사가 성격을 가른다. Re_Fin(리파이낸싱) 시나리오 상세도 같이 걸러진다.
FM_SHEET_PERIODIC = re.compile(r'\((?:q|y|fy)\)|re[_-]?fin', re.I)


def parse_xlsx(file_path, sheet_filter=None):
    """
    Excel 파일 파싱 (시트 단위 및 행 단위 병합)

    sheet_filter: 시트명 -> bool. None이면 전 시트. xlsm(재무모델)에만 건다.
    """
    sheets_data = []
    skipped = []
    try:
        wb = openpyxl.load_workbook(file_path, data_only=True)
        for sheet_name in wb.sheetnames:
            if sheet_filter and not sheet_filter(sheet_name):
                skipped.append(sheet_name)
                continue
            sheet = wb[sheet_name]
            rows = []
            row_count = 0
            chunk_idx = 1
            header = None  # 시트 첫 유효 행 = 헤더 (행 서술화용, 전략 §3.6)

            for row in sheet.iter_rows(values_only=True):
                vals = ['' if v is None else str(v).strip() for v in row]
                if not any(vals):
                    continue
                if header is None:
                    # 헤더로 쓸 만한가: 유효 칸 2개 이상
                    if sum(1 for v in vals if v) >= 2:
                        header = vals
                        rows.append(" | ".join(v for v in vals if v))
                        row_count += 1
                        continue
                    header = []  # 헤더 없는 시트로 확정

                if header and sum(1 for h in header if h) <= XLSX_PAIR_MAX_COLS:
                    # 행 서술화: '헤더=값' 쌍 — 청크가 어디서 잘려도 각 행이 자기설명적.
                    # 단 넓은 시트(분기 시계열 등)에서는 헤더 반복이 부피를 몇 배로
                    # 키우기만 하므로 값만 나열한다.
                    pairs = [f"{h}: {v}" for h, v in zip(header, vals) if v and h]
                    row_str = " | ".join(pairs) if pairs else " | ".join(v for v in vals if v)
                else:
                    row_str = " | ".join(v for v in vals if v)
                if not row_str:
                    continue
                rows.append(row_str)
                row_count += 1
                
                # 30행 단위로 청크 처리
                if row_count >= 30:
                    sheets_data.append({
                        'text': f"시트: {sheet_name}\n" + "\n".join(rows),
                        'page_number': chunk_idx,
                        'section_title': f"{sheet_name} - {chunk_idx}부",
                        'sheet_name': sheet_name,
                        'cell_range': f"1-{row_count}행"
                    })
                    rows = []
                    row_count = 0
                    chunk_idx += 1
            
            # 남은 행 처리
            if rows:
                sheets_data.append({
                    'text': f"시트: {sheet_name}\n" + "\n".join(rows),
                    'page_number': chunk_idx,
                    'section_title': f"{sheet_name} - 최종부",
                    'sheet_name': sheet_name,
                    'cell_range': f"최종 {len(rows)}개 행"
                })
    except Exception as e:
        logger.error(f"Excel 파싱 실패 ({file_path}): {e}")
    if skipped:
        logger.info(f"[xlsm] 상세 시트 {len(skipped)}개 제외 ({os.path.basename(file_path)}): "
                    + ", ".join(skipped[:8]) + ("..." if len(skipped) > 8 else ""))
    return sheets_data

def parse_doc(file_path):
    """레거시 Word(.doc, Word 97-2003) 파싱 — antiword 우선, catdoc 폴백.

    python-docx는 OOXML(.docx)만 읽는다. 코퍼스에는 확장자만 .docx로 바뀐 채
    실체는 .doc인 파일이 4건 있었고(기술규격서 3 + FS보고서), 전부 게이트에서
    'format-mismatch(ole2)'로 차단돼 통째로 빠져 있었다.
    antiword -w 0: 줄바꿈 폭 제한 해제. 표는 '|' 구분으로 렌더링된다.
    """
    import subprocess

    for cmd in (['antiword', '-w', '0', file_path], ['catdoc', file_path]):
        try:
            r = subprocess.run(cmd, capture_output=True, text=True, timeout=180)
        except FileNotFoundError:
            continue
        except Exception as e:
            logger.warning(f"{cmd[0]} 실행 실패: {e}")
            continue
        text = (r.stdout or '').strip()
        if r.returncode == 0 and len(text) >= 200:
            # antiword는 빈 줄을 많이 남긴다 — 청크 부피만 키우므로 압축한다
            text = re.sub(r'[ \t]+\n', '\n', text)
            text = re.sub(r'\n{3,}', '\n\n', text)
            logger.info(f"[{cmd[0]}] {len(text)}자 추출 ({os.path.basename(file_path)})")
            return [{'text': text, 'page_number': 1, 'section_title': '본문',
                     'sheet_name': '', 'cell_range': ''}]
    logger.warning(f"레거시 doc 추출 실패(antiword·catdoc 모두): {os.path.basename(file_path)}")
    return []


_XHTML = '{http://www.w3.org/1999/xhtml}'


def _parse_hwp_via_html(file_path, timeout=180):
    """hwp5html로 변환해 본문 + 표를 함께 뽑는다. 실패 시 빈 목록.

    **hwp5txt는 표 내용을 `<표>` 다섯 글자로 대체해 버린다.**
    실측(코퍼스 hwp/hwpx 78건): 59건(75%)이 표를 포함하고, 그 안의 값이 전부
    사라지고 있었다. '계약서_EPC가격사항'은 표 45개에 본문 3,719자 — EPC 단가가
    통째로 표 안이라 사실상 빈 문서였다. hwp5html은 같은 표를 <table>로 주므로
    pdf·docx와 동일하게 Markdown 직렬화해 되살린다.
    """
    import subprocess, tempfile, shutil
    import xml.etree.ElementTree as ET

    tmp = tempfile.mkdtemp(prefix='hwp5html_')
    try:
        r = subprocess.run(['hwp5html', '--output', tmp, file_path],
                           capture_output=True, text=True, timeout=timeout)
        index = os.path.join(tmp, 'index.xhtml')
        if r.returncode != 0 or not os.path.exists(index):
            return []
        root = ET.parse(index).getroot()

        tables = []
        for tbl in list(root.iter(_XHTML + 'table')):
            rows = []
            for tr in tbl.iter(_XHTML + 'tr'):
                cells = [' '.join(''.join(td.itertext()).split())
                         for td in tr.iter(_XHTML + 'td')]
                if any(cells):
                    rows.append(cells)
            md_parts = _rows_to_markdown(rows)
            if md_parts:
                tables.extend(md_parts)
            else:
                flat = [' | '.join(c for c in r_ if c) for r_ in rows]
                flat = [f for f in flat if f]
                if flat:
                    tables.append('\n'.join(flat))
            # 본문에서 표를 걷어내 중복 적재를 막는다. tail(표 뒤 텍스트)은 보존.
            tail = tbl.tail
            tbl.clear()
            tbl.tail = tail

        body = '\n'.join(l for l in (''.join(root.itertext())).split('\n') if l.strip())
        items = []
        if len(body.strip()) >= 20:
            items.append({'text': body.strip(), 'page_number': 1,
                          'section_title': '본문', 'sheet_name': '', 'cell_range': ''})
        for i, md in enumerate(tables):
            items.append({'text': md, 'page_number': 1, 'section_title': f'표 {i + 1}',
                          'kind': 'table', 'sheet_name': '', 'cell_range': ''})
        if items:
            logger.info(f"[hwp5html] 본문 {len(body)}자 · 표 {len(tables)}개 "
                        f"({os.path.basename(file_path)})")
        return items
    except FileNotFoundError:
        logger.warning("hwp5html 미설치 — hwp5txt로 진행 (pip install pyhwp)")
        return []
    except Exception as e:
        logger.warning(f"hwp5html 실패, hwp5txt로 진행: {e}")
        return []
    finally:
        shutil.rmtree(tmp, ignore_errors=True)


def parse_hwp(file_path):
    """한글(HWP 5.x) 파싱 — hwp5html(표 보존) → hwp5txt → 바이너리 정규식 순.

    실측(2026-07-17, 코퍼스 hwp 72건): 정규식 방식은 15건 추출 실패 + 성공분도
    텍스트를 4~10배 유실. hwp5txt는 실패 15건 중 14건 구제. (전략 §3.3, §7-5)
    2026-08-18: hwp5txt가 표를 버리는 것을 확인해 hwp5html을 1순위로 올렸다.
    """
    import subprocess

    items = _parse_hwp_via_html(file_path)
    if sum(len(i.get('text', '')) for i in items) >= 200:
        return items

    try:
        r = subprocess.run(['hwp5txt', file_path], capture_output=True, text=True, timeout=120)
        text = (r.stdout or '').strip()
        if r.returncode == 0 and len(text) >= 200:
            # hwp는 페이지 정보가 없으므로 단일 항목으로 반환한다 (분할은 청커 몫)
            return [{
                'text': text,
                'page_number': 1,
                'section_title': '본문',
                'sheet_name': '',
                'cell_range': ''
            }]
        logger.warning(f"hwp5txt 추출 부족({len(text)}자), 정규식 폴백: {os.path.basename(file_path)}")
    except FileNotFoundError:
        logger.warning("hwp5txt 미설치 — 정규식 폴백 사용 (pip install pyhwp)")
    except Exception as e:
        logger.warning(f"hwp5txt 실행 실패, 정규식 폴백: {e}")

    text_data = []
    try:
        # HWP 파일 형식은 복잡한 복합 파일 바이너리 구조입니다.
        # 정규식을 이용해 텍스트 영역을 강제로 긁어오는 fallback을 적용합니다.
        with open(file_path, 'rb') as f:
            content = f.read()
            
        # 디코딩 시도 (utf-8, cp949) 및 특수 문자 제거
        decoded_text = ""
        for encoding in ('utf-8', 'cp949', 'utf-16'):
            try:
                decoded_text = content.decode(encoding, errors='ignore')
                break
            except Exception:
                continue
        
        # 한국어 및 중요 기호 위주로 텍스트 필터링
        hangul_patt = re.compile(r'[가-힣0-9a-zA-Z\s\(\)\,\.\-\_\:\/]{2,}')
        matches = hangul_patt.findall(decoded_text)
        
        # 너무 짧은 노이즈 제거 및 단락 조립
        filtered_lines = []
        for line in matches:
            line_clean = line.strip()
            if len(line_clean) > 10 and not line_clean.startswith('SYSTEM'):
                filtered_lines.append(line_clean)
        
        # 400자 단위로 페이지를 나누어 청크화
        text_block = []
        block_idx = 1
        curr_len = 0
        for line in filtered_lines:
            text_block.append(line)
            curr_len += len(line)
            if curr_len >= 500:
                text_data.append({
                    'text': "\n".join(text_block),
                    'page_number': block_idx,
                    'section_title': f"본문 {block_idx}부",
                    'sheet_name': '',
                    'cell_range': ''
                })
                text_block = []
                curr_len = 0
                block_idx += 1
                
        if text_block:
            text_data.append({
                'text': "\n".join(text_block),
                'page_number': block_idx,
                'section_title': f"본문 {block_idx}부",
                'sheet_name': '',
                'cell_range': ''
            })
    except Exception as e:
        logger.error(f"HWP 파싱 실패 ({file_path}): {e}")
    return text_data

def parse_hwpx(file_path):
    """
    HWPX (XML 기반) 파싱
    zip 파일 내의 XML 섹션별로 텍스트를 추출하여 페이지 단위로 반환합니다.
    """
    import zipfile
    import xml.etree.ElementTree as ET
    pages = []
    try:
        with zipfile.ZipFile(file_path, 'r') as zf:
            section_files = [f for f in zf.namelist() if f.startswith('Contents/section') and f.endswith('.xml')]
            section_files.sort()
            for idx, sec_file in enumerate(section_files):
                text_block = []
                with zf.open(sec_file) as xml_file:
                    tree = ET.parse(xml_file)
                    for elem in tree.getroot().iter():
                        if elem.tag.endswith('}t') and elem.text:
                            text_block.append(elem.text)
                if text_block:
                    pages.append({
                        'text': "\n".join(text_block),
                        'page_number': idx + 1,
                        'section_title': f"섹션 {idx + 1}",
                        'sheet_name': '',
                        'cell_range': ''
                    })
    except Exception as e:
        logger.error(f"HWPX 파싱 실패 ({file_path}): {e}")
    return pages

def parse_pptx(file_path):
    """
    PowerPoint 파일 파싱 (슬라이드 단위로 텍스트 획득)
    """
    from pptx import Presentation
    slides_data = []
    try:
        prs = Presentation(file_path)
        for idx, slide in enumerate(prs.slides):
            slide_text = []
            slide_tables = []

            # 슬라이드 내 모든 도형(shape)에서 텍스트·표 수집
            # (표 도형은 .text가 없어 기존 코드에서 통째로 유실되던 갭 수정 — 전략 §3.6)
            for shape in slide.shapes:
                if getattr(shape, 'has_table', False):
                    rows = [[cell.text for cell in row.cells] for row in shape.table.rows]
                    md_parts = _rows_to_markdown(rows)
                    if md_parts:
                        slide_tables.extend(md_parts)
                    else:
                        flat = [" | ".join(c.strip() for c in r if c.strip()) for r in rows]
                        slide_text.extend(f for f in flat if f)
                elif hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())

            if slide_text:
                joined_text = "\n".join(slide_text)
                slides_data.append({
                    'text': joined_text,
                    'page_number': idx + 1,
                    'section_title': f"슬라이드 {idx + 1}",
                    'sheet_name': '',
                    'cell_range': ''
                })
            for md in slide_tables:
                slides_data.append({
                    'text': md,
                    'page_number': idx + 1,
                    'section_title': f"표 (슬라이드 {idx + 1})",
                    'kind': 'table',
                    'sheet_name': '',
                    'cell_range': ''
                })
    except Exception as e:
        logger.error(f"PPTX 파싱 실패 ({file_path}): {e}")
    return slides_data

def parse_markdown(file_path):
    """마크다운/텍스트 파싱 — 헤딩(#~######) 단위로 섹션을 나눈다.

    사업개요처럼 사람이 직접 관리하는 정리본을 위한 경로다.
    헤딩을 그대로 section_title로 남기는 것이 핵심이다. 사업개요 양식은
    `## 4. 인허가 현황 — 당진행복솔라`처럼 **섹션 제목에 사업명을 반복**하도록
    설계돼 있어, 이 값이 출처 표시와 검색 맥락을 동시에 책임진다.

    표는 별도 항목(kind='table')으로 떼지 않고 섹션 안에 그대로 둔다.
    떼어내면 "어느 사업 표인지" 알 수 없는 조각이 되기 때문이다.
    """
    text = ''
    for enc in ('utf-8', 'utf-8-sig', 'cp949'):
        try:
            with open(file_path, 'r', encoding=enc) as f:
                text = f.read()
            break
        except UnicodeDecodeError:
            continue
    if not text.strip():
        logger.warning(f"마크다운 본문이 비어 있습니다: {file_path}")
        return []

    heading_re = re.compile(r'^(#{1,6})\s+(.+?)\s*$', re.MULTILINE)
    matches = list(heading_re.finditer(text))

    items, idx = [], 1

    def _add(title, body):
        nonlocal idx
        body = (body or '').strip()
        if not body:
            return
        items.append({
            'text': body,
            'page_number': idx,       # md에는 페이지가 없다. 섹션 순번을 페이지로 쓴다
            'section_title': title,
            'sheet_name': '',
            'cell_range': '',
        })
        idx += 1

    if not matches:
        _add(os.path.splitext(os.path.basename(file_path))[0], text)
        return items

    # 첫 헤딩 앞의 머리말(제목·메타 블록)
    _add('머리말', text[:matches[0].start()])

    for i, m in enumerate(matches):
        title = m.group(2).strip()
        end = matches[i + 1].start() if i + 1 < len(matches) else len(text)
        # 헤딩 줄을 본문에 포함시킨다 — 조각만 검색돼도 무슨 섹션인지 알 수 있다
        _add(title, text[m.start():end])

    return items


def parse_file(file_path):
    """
    파일 확장자에 맞춰 해당 파서 함수를 호출하는 범용 엔드포인트.

    1) 매직바이트 게이트 — 위장/DRM 파일을 사유와 함께 차단 (ParserGateError)
    2) 포맷별 파싱 (표는 kind='table' 항목으로 분리되어 나옴)
    3) 머리말/꼬리말 제거 — 본문 항목에만 적용 (표/OCR 항목은 보호)
    """
    ext = os.path.splitext(file_path)[1].lower().strip('.')
    gate_opts = gate_check(file_path, ext)  # 불가 시 ParserGateError 전파
    ext = gate_opts.get('as_ext', ext)      # 확장자 위장 교정(hwpx→hwp 등)

    if ext == 'pdf':
        items = parse_pdf(file_path, pdf_offset=gate_opts.get('pdf_offset', 0))
    elif ext == 'docx':
        items = parse_docx(file_path)
    elif ext == 'doc':
        items = parse_doc(file_path)
    elif ext in ('xlsx', 'xlsm'):
        # xlsm은 매크로가 든 통합문서일 뿐 시트 구조는 xlsx와 같다.
        # openpyxl은 셀 값만 읽고 VBA는 실행하지 않으므로 파싱 자체는 안전하다.
        #
        # 화이트리스트는 xlsm에만 건다. 코퍼스의 xlsm은 전부 재무모델이고,
        # xlsx(자금집행현황·인출요청서 등)는 운영 실무 자료라 시트를 걸러내면 안 된다.
        sf = (lambda n: bool(FM_SHEET_WHITELIST.search(n))
                        and not FM_SHEET_PERIODIC.search(n)) if ext == 'xlsm' else None
        items = parse_xlsx(file_path, sheet_filter=sf)
    elif ext == 'hwp':
        items = parse_hwp(file_path)
    elif ext == 'hwpx':
        items = parse_hwpx(file_path)
    elif ext == 'pptx':
        items = parse_pptx(file_path)
    elif ext in ('md', 'markdown', 'txt'):
        items = parse_markdown(file_path)
    else:
        logger.warning(f"지원되지 않는 확장자 ({ext}): {file_path}")
        return []

    items = items or []

    # NUL(0x00) 제거 — Postgres는 문자열에 NUL을 허용하지 않는다.
    # (실측: 일부 PDF의 임베디드 폰트 추출물에 NUL 포함 → DocumentChunk 저장 실패)
    for item in items:
        if item.get('text') and '\x00' in item['text']:
            item['text'] = item['text'].replace('\x00', '')

    # hwp 정규식 파서는 추출 실패 시 잡음만 남는다 — 200자 미만이면 실패로 명시 (전략 §3.3)
    if ext == 'hwp':
        total = sum(len(i.get('text', '')) for i in items)
        if total < 200:
            raise ParserGateError(f'hwp-extract-too-short({total})')

    # Excel은 시트/셀 단위라 '페이지 가장자리' 개념이 없으므로 제외한다.
    if ext in ('xlsx', 'xlsm'):
        return items

    # 표·OCR 항목은 머리말 제거 대상에서 보호 (Markdown 표의 첫 줄은 헤더다)
    body = [i for i in items if i.get('kind') != 'table']
    tables = [i for i in items if i.get('kind') == 'table']
    return strip_repeated_edges(body) + tables

